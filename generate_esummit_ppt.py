"""
VidyaMitra — E-Summit Hackathon Submission PPT Generator
Generates a professional, dark-themed PPTX presentation following the EXACT 7-slide format.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
BG_ACCENT = RGBColor(0x25, 0x25, 0x3A)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xA0, 0xB8)
MUTED = RGBColor(0x72, 0x72, 0x8F)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_bg(slide, left, top, width, height, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=LIGHT_GRAY, bullet_color=CYAN, spacing=Pt(10), bold_first_word=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing

        if bold_first_word and ":" in item:
            parts = item.split(":", 1)
            run1 = p.add_run()
            run1.text = parts[0] + ":"
            run1.font.bold = True
            run1.font.color.rgb = WHITE
            
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.bold = False
        else:
            run = p.add_run()
            run.text = item
            
    return txBox

def create_header(slide, title):
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                 title, font_size=32, color=WHITE, bold=True)
    add_shape_bg(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.02), BG_CARD)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE (Format Exact Match)
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "SMART CAMPUS HACKATHON 2026", font_size=44, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.8),
             "VidyaMitra: Smart Skill-to-Career Mapping & Opportunity\nRecommendation Platform",
             font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Info box
add_shape_bg(slide, Inches(3.2), Inches(4.0), Inches(7.0), Inches(2.8), BG_CARD, radius=0.05)

info_lines = [
    "Team Name: [Your Team Name]",
    "Problem Statement Title: Smart Skill-to-Career Mapping & Opportunity Recommendation Platform for Engineering Students",
    "Team Members: [Member 1], [Member 2], [Member 3], [Member 4]",
    "College / Organization: [Your College Name]",
    "MAIL ID: [Your Email]"
]

for i, line in enumerate(info_lines):
    parts = line.split(":", 1)
    
    add_text_box(slide, Inches(3.5), Inches(4.2 + i*0.45), Inches(2.5), Inches(0.4),
                 parts[0] + ":", font_size=14, color=CYAN, bold=True)
                 
    if len(parts) > 1:
        add_text_box(slide, Inches(6.0), Inches(4.2 + i*0.45), Inches(4.0), Inches(0.4),
                     parts[1].strip(), font_size=14, color=WHITE, bold=False)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
create_header(slide, "PROBLEM STATEMENT")

# Clearly define the problem you are addressing
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.5), BG_CARD, radius=0.03)
add_text_box(slide, Inches(1.1), Inches(1.8), Inches(11), Inches(0.4),
             "🎯 The Core Problem", font_size=18, color=CYAN, bold=True)
add_text_box(slide, Inches(1.1), Inches(2.3), Inches(11), Inches(0.6),
             "Engineering students gain academic knowledge but lack visibility into evolving industry-required skills. They struggle to identify suitable opportunities (internships, hackathons, certifications) because these are scattered across multiple platforms. Moreover, students are unsure if their resumes meet ATS eligibility criteria, leading to misaligned applications and missed opportunities.",
             font_size=15, color=LIGHT_GRAY)

# Target users
add_shape_bg(slide, Inches(0.8), Inches(3.3), Inches(5.7), Inches(3.5), BG_CARD, radius=0.03)
add_text_box(slide, Inches(1.1), Inches(3.5), Inches(5), Inches(0.4),
             "👥 Target Users & Stakeholders", font_size=18, color=PURPLE, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(4.1), Inches(5.1), Inches(2.5), [
    "Engineering Students: Across all branches (CSE, ECE, EEE, Mech, Civil) looking for jobs, internships, or skills.",
    "Placement Teams / Universities: Wanting cohort-level analytics to improve batch placement readiness.",
    "Recruiters & Companies: Looking for skill-verified candidates matched to specific JD requirements."
], font_size=14, spacing=Pt(14))

# Why important
add_shape_bg(slide, Inches(6.8), Inches(3.3), Inches(5.7), Inches(3.5), BG_CARD, radius=0.03)
add_text_box(slide, Inches(7.1), Inches(3.5), Inches(5), Inches(0.4),
             "❗ Why This Problem is Important", font_size=18, color=ROSE, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(4.1), Inches(5.1), Inches(2.5), [
    "The Skills Gap: 80% of graduates are deemed 'unemployable' due to a mismatch between academia and industry needs.",
    "Inefficient Discovery: Without personalized mapping, students apply blindly and face high rejection rates.",
    "Data-Blind Institutions: Colleges lack real-time visibility into what skills their batch is missing until after placement season fails."
], font_size=14, spacing=Pt(14))


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: BACKGROUND AND MOTIVATION
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
create_header(slide, "BACKGROUND AND MOTIVATION")

# Current Scenario / Limitations
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2), BG_CARD, radius=0.03)
add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.4),
             "📉 Current Scenario & Limitations", font_size=20, color=CYAN, bold=True)

add_bullet_list(slide, Inches(1.1), Inches(2.5), Inches(5.1), Inches(4.0), [
    "Static Skill Tracking: Current systems track grades (CGPA), not specific granular skills mapped to industry roles.",
    "Fragmented Ecosystem: Job boards (internships), Devfolio (hackathons), and Coursera (certifications) exist independently without cross-referencing a student's profile.",
    "Generic Advice: Most career portals only cater to Software/CS engineers, leaving ECE, Mechanical, and Civil students with generic, unactionable advice.",
    "Reactive Intervention: Colleges act on placement failures, rather than proactively monitoring student 'readiness scores' through data."
], font_size=15, spacing=Pt(20))

# Motivation
add_shape_bg(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2), BG_CARD, radius=0.03)
add_text_box(slide, Inches(7.1), Inches(1.8), Inches(5), Inches(0.4),
             "🚀 Our Motivation", font_size=20, color=PURPLE, bold=True)

add_text_box(slide, Inches(7.1), Inches(2.5), Inches(5.1), Inches(1.0),
             "We wanted to build an 'Intelligent Career OS' that serves as a constant companion (VidyaMitra) throughout a student's degree. Not just a job board, but a personalized readiness engine.",
             font_size=15, color=WHITE)

add_bullet_list(slide, Inches(7.1), Inches(3.8), Inches(5.1), Inches(2.8), [
    "Democratize Opportunity Recommendations: Match the right hackathon or certification to the exact missing skill string.",
    "Bring Industry Benchmarks to Campus: Show students exactly what percentage of JDs demand a certain skill today.",
    "Empower Universities: Give placement officers 'God-mode' analytics over their cohort's real-time skill heatmaps."
], font_size=15, spacing=Pt(16))


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: PROPOSED SOLUTION (WITH FLOWCHART)
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
create_header(slide, "PROPOSED SOLUTION (WITH FLOWCHART)")

# Text section (Left)
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), BG_CARD, radius=0.03)
add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.4),
             "✨ Key Features & Capabilities", font_size=20, color=CYAN, bold=True)

add_bullet_list(slide, Inches(1.1), Inches(2.5), Inches(5.0), Inches(4.1), [
    "ATS Resume Score Engine: Upload PDF, get instant compatibility score, keyword matching, formatting review.",
    "Skill Gap & Benchmark Analyzer: Weighted matching linking 25+ roles to mapped skills over 1200+ actual JDs.",
    "Curated Recommendations: Hackathons (SIH, MLH) & Certs (AWS, NPTEL) suggested based ONLY on missing skills.",
    "AI Learning Roadmap: 54KB engine creates multi-phase weekly plans.",
    "Placement Dashboard: Heatmaps and readiness metrics for college admin to identify at-risk students."
], font_size=13, spacing=Pt(12))

# Flowchart Section (Right)
add_shape_bg(slide, Inches(6.6), Inches(1.6), Inches(5.9), Inches(5.2), BG_CARD, radius=0.03)
add_text_box(slide, Inches(6.9), Inches(1.8), Inches(5), Inches(0.4),
             "🔄 Solution Flowchart", font_size=20, color=PURPLE, bold=True)

# Build a visual flowchart using shapes
steps = [
    ("Student Profile\n(Resume + Domain)", Inches(7.3), Inches(2.4), CYAN),
    ("↓  AI Extraction & ATS Scoring  ↓", Inches(7.3), Inches(3.2), MUTED),
    ("Skill Gap Analysis\nvs Industry Benchmarks", Inches(7.3), Inches(3.7), PURPLE),
    ("↙                                  ↘", Inches(7.3), Inches(4.5), MUTED),
    ("Personalized Recs:\nHackathons & Certs", Inches(6.8), Inches(5.0), EMERALD),
    ("AI Generates\nLearning Roadmap", Inches(9.8), Inches(5.0), AMBER),
    ("↓  Feeds Data to Placement Team  ↓", Inches(7.3), Inches(5.9), MUTED),
    ("Cohort Analytics Dashboard", Inches(7.3), Inches(6.3), ROSE),
]

for text, x, y, color in steps:
    if "↓" in text or "↙" in text:
        add_text_box(slide, x, y, Inches(4.5), Inches(0.4), text, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    else:
        width = Inches(2.5) if "Recs" in text or "Roadmap" in text else Inches(4.5)
        # add_shape_bg(slide, x, y, width, Inches(0.8), color, radius=0.1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: TECHNOLOGY DETAILS
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
create_header(slide, "TECHNOLOGY DETAILS")

# Columns for tech stack
tech_stacks = [
    ("Frontend Layer", "🖥️", [
        ("React 19 & Vite 8", "High-performance SPA"),
        ("Framer Motion", "Complex micro-animations"),
        ("Recharts", "Interactive analytics graphs"),
        ("Glassmorphism", "Premium custom UI/UX design")
    ], CYAN),
    ("Backend & Database", "⚙️", [
        ("Python FastAPI", "Asynchronous API logic (17 routers)"),
        ("Supabase (PostgreSQL)", "DB, Auth, and Row Level Security"),
        ("Llama 4 Scout (Groq)", "Ultra-fast LLM inference for AI"),
        ("PyMuPDF", "Resume text extraction")
    ], PURPLE),
    ("AI / Logic Engines", "🧠", [
        ("Skill Gap Engine", "Deterministic weighted skill matching"),
        ("Recommendation System", "Profile to opportunity mapping"),
        ("Roadmap Engine", "Multi-phase schedule generation"),
        ("Placement Engine", "Cohort data aggregation heuristics")
    ], EMERALD),
]

for i, (title, emoji, items, color) in enumerate(tech_stacks):
    x = Inches(0.8 + i*4.0)
    add_shape_bg(slide, x, Inches(1.6), Inches(3.7), Inches(3.8), BG_CARD, radius=0.03)
    add_shape_bg(slide, x, Inches(1.6), Inches(3.7), Inches(0.05), color)
    
    add_text_box(slide, x + Inches(0.2), Inches(1.8), Inches(3.3), Inches(0.4),
                 f"{emoji} {title}", font_size=18, color=color, bold=True)
                 
    bullet_list = [f"{tech}: {desc}" for tech, desc in items]
    add_bullet_list(slide, x + Inches(0.2), Inches(2.4), Inches(3.3), Inches(2.8),
                    bullet_list, font_size=12, spacing=Pt(12))

# Architecture Description & Justification
add_shape_bg(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.3), BG_CARD, radius=0.03)
add_text_box(slide, Inches(1.1), Inches(5.8), Inches(11.1), Inches(0.3),
             "Basic Architecture & Technology Justification", font_size=16, color=AMBER, bold=True)
add_text_box(slide, Inches(1.1), Inches(6.1), Inches(11.1), Inches(0.6),
             "Our architecture separates deterministic matching (reliable logic) from AI generation (LLM tasks) to prevent hallucination. We use FastAPI for rapid parallel execution and Groq API for near-instant (LPU) AI response times, crucial for interactive student tools. Supabase provides out-of-the-box secure role segregation between Students and Placement Officers.",
             font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: IMPACT AND OUTCOMES
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
create_header(slide, "IMPACT AND OUTCOMES")

add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.2), BG_CARD, radius=0.03)
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(3.7), Inches(0.05), EMERALD)
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(3.3), Inches(0.4),
             "🌟 Expected Impact", font_size=20, color=EMERALD, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.5), Inches(3.3), Inches(4.0), [
    "Higher Placement Rates: Students target exactly what industry wants.",
    "Reduced ATS Rejection: Resumes are pre-verified and tailored before application.",
    "Proactive Institutions: Colleges intervene before Day 0 placements.",
    "Opportunity Discovery: Hackathon / Certification participation triples due to contextual surfacing."
], font_size=14, spacing=Pt(18))

add_shape_bg(slide, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.2), BG_CARD, radius=0.03)
add_shape_bg(slide, Inches(4.8), Inches(1.6), Inches(3.7), Inches(0.05), CYAN)
add_text_box(slide, Inches(5.0), Inches(1.8), Inches(3.3), Inches(0.4),
             "💡 Innovative Aspects", font_size=20, color=CYAN, bold=True)
add_bullet_list(slide, Inches(5.0), Inches(2.5), Inches(3.3), Inches(4.0), [
    "Multi-Branch AI: Supports ECE/EEE/Mech/Civil (25+ roles) not generic CS.",
    "Live Benchmarking data framing skills vs real JD percentages.",
    "Dual-Faced Dashboard: Same data renders as Personal roadmaps for students, but Cohort Heatmaps for officers.",
    "Skill-Matched Hackathons over manual searching."
], font_size=14, spacing=Pt(18))

add_shape_bg(slide, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.2), BG_CARD, radius=0.03)
add_shape_bg(slide, Inches(8.8), Inches(1.6), Inches(3.7), Inches(0.05), PURPLE)
add_text_box(slide, Inches(9.0), Inches(1.8), Inches(3.3), Inches(0.4),
             "📈 Scalability & Future", font_size=20, color=PURPLE, bold=True)
add_bullet_list(slide, Inches(9.0), Inches(2.5), Inches(3.3), Inches(4.0), [
    "Live JD Scraping: Auto-updating skill requirements in real-time.",
    "B2B SaaS Model: License to universities as 'Placement OS'.",
    "LMS Integration: Sync with Moodle/Blackboard to verify academic skills.",
    "Alumni Network: Match students with alumni currently in target roles."
], font_size=14, spacing=Pt(18))


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: REFERENCES
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
create_header(slide, "REFERENCES")

add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2), BG_CARD, radius=0.03)

refs = [
    "1. React 19 Documentation & Best Practices (https://react.dev/)",
    "2. Node.js & Vite Build Optimizations (https://vitejs.dev/)",
    "3. Supabase / PostgreSQL Architecture (https://supabase.com/docs)",
    "4. FastAPI Asynchronous Programming (https://fastapi.tiangolo.com/)",
    "5. Groq LPU API / LLaVa / Llama 4 Model Documentation (https://console.groq.com/docs/models)",
    "6. AI in EdTech: Addressing the Skills Gap (Industry Journal Articles)",
    "7. Applicant Tracking System (ATS) Keyword Heuristics & Parsing Patterns",
    "8. Job Description Skills Frequency Distribution Methods (Kaggle Datasets)"
]

add_bullet_list(slide, Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.0), 
               refs, font_size=16, spacing=Pt(18), bold_first_word=False)

# Add Thank You at bottom
add_shape_bg(slide, Inches(4.0), Inches(6.0), Inches(5.333), Inches(0.8), rgba_to_hex(PURPLE) if False else PURPLE, radius=0.1)
add_text_box(slide, Inches(4.0), Inches(6.2), Inches(5.333), Inches(0.8),
             "Thank You! Let's bridge the skill gap.", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════

output_path = "/Users/srimannarayanadeevi/Documents/vidyamitra/VidyaMitra_ESummit_PPT.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
