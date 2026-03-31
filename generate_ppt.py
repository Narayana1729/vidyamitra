"""
VidyaMitra — Hackathon Submission PPT Generator
Generates a professional, dark-themed PPTX presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
BG_ACCENT = RGBColor(0x25, 0x25, 0x3A)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xA0, 0xB8)
MUTED = RGBColor(0x72, 0x72, 0x8F)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_bg(slide, left, top, width, height, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=LIGHT_GRAY, bullet_color=CYAN, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Accent bar
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
             "VidyaMitra", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "AI-Powered Career Intelligence Platform for Engineering Students",
             font_size=24, color=CYAN, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.8),
             "Smart Skill-to-Career Mapping & Opportunity Recommendation",
             font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

# Divider
add_shape_bg(slide, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.04), PURPLE)

add_text_box(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.6),
             "Problem Statement #1  •  Hackathon Submission 2026",
             font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.6),
             "Team: [Your Names Here]",
             font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "THE PROBLEM", font_size=14, color=PURPLE, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.8),
             "Engineering students lack a unified platform for career readiness",
             font_size=32, color=WHITE, bold=True)

# Problem cards
problems = [
    ("🎯", "Scattered Opportunities", "Internships, hackathons, certifications scattered across 10+ platforms. Students miss relevant opportunities."),
    ("📊", "Invisible Skill Gaps", "No structured system to evaluate skills against industry expectations. Students don't know what they're missing."),
    ("📄", "Resume Uncertainty", "Students unsure if their resumes meet ATS criteria. Result: rejected before a human reads it."),
    ("🏫", "Placement Teams Blind", "No cohort-level analytics. Placement officers can't identify at-risk students or department-wide gaps."),
]

for i, (emoji, title, desc) in enumerate(problems):
    col = i % 4
    x = Inches(0.6 + col * 3.1)
    y = Inches(2.3)

    add_shape_bg(slide, x, y, Inches(2.9), Inches(3.8), BG_CARD, radius=0.05)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.3), Inches(0.5),
                 emoji, font_size=36, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.3), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.6), Inches(2.3), Inches(2.0),
                 desc, font_size=12, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
             "Result: Misaligned applications, missed opportunities, and ineffective career preparation.",
             font_size=14, color=ROSE, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: OUR SOLUTION
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "OUR SOLUTION", font_size=14, color=PURPLE, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.8),
             "VidyaMitra — One platform. Complete career lifecycle.",
             font_size=32, color=WHITE, bold=True)

features = [
    ("📄", "Resume Analysis", "ATS scoring, keyword\nmatch, AI suggestions", CYAN),
    ("🧠", "Skill Gap Analysis", "25+ roles × 5 branches\nweighted matching", PURPLE),
    ("📊", "Industry Benchmarks", "1,200+ JDs analyzed from\nLinkedIn, Naukri, Indeed", AMBER),
    ("🏆", "Hackathon Recs", "15 hackathons matched\nto your skill gaps", EMERALD),
    ("📜", "Certification Recs", "20 certifications from\nAWS, NPTEL, Google", CYAN),
    ("🗺️", "Learning Roadmap", "AI-generated phased\nlearning pathways", PURPLE),
    ("🎤", "Mock Interview", "AI questions + live\nanswer evaluation", AMBER),
    ("📊", "Placement Dashboard", "Cohort readiness, dept\nheatmap, at-risk alerts", ROSE),
]

for i, (emoji, title, desc, color) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.1)
    y = Inches(2.2 + row * 2.4)

    add_shape_bg(slide, x, y, Inches(2.9), Inches(2.1), BG_CARD, radius=0.05)

    # Color accent bar
    add_shape_bg(slide, x, y, Inches(2.9), Inches(0.05), color)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.4),
                 emoji, font_size=22)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.2), Inches(2.0), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(2.3), Inches(1.2),
                 desc, font_size=11, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: PS MAPPING TABLE
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "PROBLEM STATEMENT COVERAGE", font_size=14, color=PURPLE, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.8),
             "Every requirement mapped. 93% coverage.",
             font_size=32, color=WHITE, bold=True)

mappings = [
    ("\"Analyzes student profiles\"", "Resume Analyzer + Skill Profiling + Domain Detection", "✅"),
    ("\"Visualizes skill gaps\"", "Score gauge + priority badges + demand bars", "✅"),
    ("\"Assesses resume readiness\"", "ATS Score Engine — keyword + section analysis", "✅"),
    ("\"Internship recommendations\"", "Job Board with company postings + match scoring", "✅"),
    ("\"Project recommendations\"", "AI-generated domain-aware project suggestions", "✅"),
    ("\"Hackathon recommendations\"", "15 hackathons matched by skills + domain", "✅"),
    ("\"Certification recommendations\"", "20 certifications (AWS, NPTEL, Google, SolidWorks)", "✅"),
    ("\"Learning pathways\"", "Multi-phase AI roadmap with weekly schedule", "✅"),
    ("\"Dashboards for students\"", "Readiness gauge, charts, insights, activity feed", "✅"),
    ("\"Analytics for placement teams\"", "Cohort heatmap, dept readiness, at-risk students", "✅"),
]

# Table header
add_shape_bg(slide, Inches(0.6), Inches(2.0), Inches(5.0), Inches(0.4), PURPLE)
add_text_box(slide, Inches(0.8), Inches(2.03), Inches(4.5), Inches(0.4),
             "Problem Statement Says", font_size=11, color=WHITE, bold=True)

add_shape_bg(slide, Inches(5.6), Inches(2.0), Inches(6.5), Inches(0.4), PURPLE)
add_text_box(slide, Inches(5.8), Inches(2.03), Inches(5.5), Inches(0.4),
             "VidyaMitra Solution", font_size=11, color=WHITE, bold=True)

add_shape_bg(slide, Inches(12.1), Inches(2.0), Inches(0.8), Inches(0.4), PURPLE)
add_text_box(slide, Inches(12.1), Inches(2.03), Inches(0.8), Inches(0.4),
             "Status", font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, (ps, solution, status) in enumerate(mappings):
    y = Inches(2.45 + i * 0.44)
    bg_color = BG_CARD if i % 2 == 0 else BG_ACCENT

    add_shape_bg(slide, Inches(0.6), y, Inches(5.0), Inches(0.42), bg_color)
    add_text_box(slide, Inches(0.8), y + Inches(0.03), Inches(4.5), Inches(0.38),
                 ps, font_size=10, color=LIGHT_GRAY)

    add_shape_bg(slide, Inches(5.6), y, Inches(6.5), Inches(0.42), bg_color)
    add_text_box(slide, Inches(5.8), y + Inches(0.03), Inches(5.5), Inches(0.38),
                 solution, font_size=10, color=WHITE)

    add_shape_bg(slide, Inches(12.1), y, Inches(0.8), Inches(0.42), bg_color)
    add_text_box(slide, Inches(12.1), y + Inches(0.03), Inches(0.8), Inches(0.38),
                 status, font_size=14, color=EMERALD, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: MULTI-BRANCH SUPPORT
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "MULTI-BRANCH ENGINEERING SUPPORT", font_size=14, color=PURPLE, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.8),
             "Not just CS — all engineering branches supported",
             font_size=32, color=WHITE, bold=True)

branches = [
    ("💻", "CSE / IT", "25+ roles", "Python, React, Docker, AWS, SQL", "1,247 JDs", PURPLE),
    ("📡", "ECE", "5 roles", "Embedded, VLSI, Verilog, MATLAB", "487 JDs", CYAN),
    ("⚡", "EEE", "4 roles", "Power Systems, PLC, MATLAB, ETAP", "312 JDs", AMBER),
    ("⚙️", "Mechanical", "5 roles", "SolidWorks, ANSYS, AutoCAD, GD&T", "623 JDs", EMERALD),
    ("🏗️", "Civil", "5 roles", "AutoCAD, Revit, STAAD.Pro, Primavera", "389 JDs", ROSE),
]

for i, (emoji, name, roles, skills, jds, color) in enumerate(branches):
    x = Inches(0.4 + i * 2.55)
    y = Inches(2.3)

    add_shape_bg(slide, x, y, Inches(2.4), Inches(4.5), BG_CARD, radius=0.05)
    add_shape_bg(slide, x, y, Inches(2.4), Inches(0.05), color)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.25), Inches(2.0), Inches(0.5),
                 f"{emoji} {name}", font_size=20, color=WHITE, bold=True)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.0), Inches(0.3),
                 roles, font_size=12, color=color, bold=True)

    add_text_box(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.0), Inches(0.3),
                 "Key Skills:", font_size=10, color=MUTED, bold=True)

    add_text_box(slide, x + Inches(0.2), y + Inches(1.7), Inches(2.0), Inches(1.5),
                 skills, font_size=11, color=LIGHT_GRAY)

    add_text_box(slide, x + Inches(0.2), y + Inches(3.5), Inches(2.0), Inches(0.3),
                 f"📊 {jds} Analyzed", font_size=11, color=color, bold=True)

add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4),
             "Domain-specific roles, skills, projects, hackathons, and certifications for each branch",
             font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: ARCHITECTURE
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "ARCHITECTURE & TECH STACK", font_size=14, color=PURPLE, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6),
             "Full-stack platform with AI-enriched engines",
             font_size=28, color=WHITE, bold=True)

# Frontend box
add_shape_bg(slide, Inches(0.6), Inches(2.0), Inches(3.8), Inches(4.8), BG_CARD, radius=0.04)
add_shape_bg(slide, Inches(0.6), Inches(2.0), Inches(3.8), Inches(0.05), CYAN)
add_text_box(slide, Inches(0.9), Inches(2.2), Inches(3.2), Inches(0.4),
             "🖥️  FRONTEND", font_size=16, color=CYAN, bold=True)
add_bullet_list(slide, Inches(0.9), Inches(2.8), Inches(3.2), Inches(3.5), [
    "React 19 + Vite 8",
    "Framer Motion animations",
    "Recharts data visualization",
    "Lucide React icons",
    "Dark / Light theme",
    "Glassmorphism design",
    "Role-based layouts",
    "LocalStorage + Supabase sync",
], font_size=12)

# Backend box
add_shape_bg(slide, Inches(4.7), Inches(2.0), Inches(4.2), Inches(4.8), BG_CARD, radius=0.04)
add_shape_bg(slide, Inches(4.7), Inches(2.0), Inches(4.2), Inches(0.05), PURPLE)
add_text_box(slide, Inches(5.0), Inches(2.2), Inches(3.6), Inches(0.4),
             "⚙️  BACKEND (FastAPI)", font_size=16, color=PURPLE, bold=True)
add_bullet_list(slide, Inches(5.0), Inches(2.8), Inches(3.6), Inches(3.5), [
    "17 API routers",
    "Skill Gap Engine (weighted matching)",
    "ATS Score Engine",
    "AI Roadmap Generator (54KB engine)",
    "Interview Simulator + Evaluator",
    "Recommendations Engine",
    "Placement Analytics Engine",
    "Llama 4 Scout via Groq API",
], font_size=12)

# Database box
add_shape_bg(slide, Inches(9.2), Inches(2.0), Inches(3.6), Inches(4.8), BG_CARD, radius=0.04)
add_shape_bg(slide, Inches(9.2), Inches(2.0), Inches(3.6), Inches(0.05), EMERALD)
add_text_box(slide, Inches(9.5), Inches(2.2), Inches(3.0), Inches(0.4),
             "🗄️  DATABASE", font_size=16, color=EMERALD, bold=True)
add_bullet_list(slide, Inches(9.5), Inches(2.8), Inches(3.0), Inches(3.5), [
    "Supabase (PostgreSQL)",
    "10 tables with FK relationships",
    "Role-based access control",
    "Real-time subscriptions",
    "Row-level security",
    "profiles, resume_analyses,",
    "skill_analyses, roadmaps,",
    "jobs, applications, activity_log",
], font_size=12)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: DEMO FLOW
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "DEMO FLOW", font_size=14, color=PURPLE, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6),
             "Live walkthrough — from resume to placement analytics",
             font_size=28, color=WHITE, bold=True)

steps = [
    ("1", "Upload Resume", "Student uploads PDF → ATS Score (72%)\nSection analysis + suggestions", CYAN),
    ("2", "Skill Gap Analysis", "Select target role → Add skills →\n65% match, prioritized gaps", PURPLE),
    ("3", "Industry Benchmarks", "\"1,247 JDs analyzed\" → demand %\nbars for top skills", AMBER),
    ("4", "Recommendations", "Matched hackathons (SIH, MLH)\n+ certifications (AWS, NPTEL)", EMERALD),
    ("5", "Learning Roadmap", "AI generates 4-phase plan\nwith weekly schedule", CYAN),
    ("6", "Mock Interview", "AI asks questions → evaluates\nanswers → gives feedback", PURPLE),
    ("7", "Placement Dashboard", "Officer sees 72% batch readiness\ndept heatmap, at-risk students", ROSE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.3 + (i % 4) * 3.2)
    y = Inches(2.1 + (i // 4) * 2.7)

    add_shape_bg(slide, x, y, Inches(3.0), Inches(2.3), BG_CARD, radius=0.05)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, x + Inches(0.2), y + Inches(0.22), Inches(0.45), Inches(0.45),
                 num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.8), y + Inches(0.25), Inches(2.0), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.85), Inches(2.4), Inches(1.2),
                 desc, font_size=11, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: USPs
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "WHAT MAKES US DIFFERENT", font_size=14, color=PURPLE, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6),
             "6 USPs that set VidyaMitra apart",
             font_size=28, color=WHITE, bold=True)

usps = [
    ("🔀", "Multi-Branch", "Works for ECE, EEE, Mech, Civil\n— not just Computer Science", CYAN),
    ("📊", "Industry-Benchmarked", "Skills mapped against 1,200+ real\nJDs, not static lists", AMBER),
    ("🔄", "End-to-End Pipeline", "Resume → Skills → Roadmap →\nInterview → Job → Track", PURPLE),
    ("👥", "Dual Dashboard", "Students: personal analytics.\nPlacement: cohort analytics", EMERALD),
    ("🧠", "AI + Deterministic", "Core matching is reliable.\nAI enhances, never replaces", ROSE),
    ("🏫", "Campus-Deployable", "Designed for colleges. Dept grouping,\nbatch analytics, officer views", CYAN),
]

for i, (emoji, title, desc, color) in enumerate(usps):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(2.1 + row * 2.6)

    add_shape_bg(slide, x, y, Inches(3.9), Inches(2.2), BG_CARD, radius=0.05)
    add_shape_bg(slide, x, y, Inches(0.06), Inches(2.2), color)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.3), Inches(0.5),
                 f"{emoji}  {title}", font_size=18, color=WHITE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.3), Inches(1.2),
                 desc, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: NUMBERS
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "BY THE NUMBERS", font_size=14, color=PURPLE, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.6),
             "The scale of VidyaMitra",
             font_size=28, color=WHITE, bold=True)

numbers = [
    ("17", "API Endpoints", PURPLE),
    ("25+", "Career Roles", CYAN),
    ("5", "Engg Branches", EMERALD),
    ("1,200+", "JDs Analyzed", AMBER),
    ("15", "Hackathons", ROSE),
    ("20", "Certifications", CYAN),
    ("10", "DB Tables", PURPLE),
    ("54KB", "Roadmap Engine", EMERALD),
]

for i, (num, label, color) in enumerate(numbers):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(2.2 + row * 2.5)

    add_shape_bg(slide, x, y, Inches(2.8), Inches(2.1), BG_CARD, radius=0.05)

    add_text_box(slide, x, y + Inches(0.3), Inches(2.8), Inches(0.8),
                 num, font_size=44, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.2), Inches(2.8), Inches(0.5),
                 label, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: FUTURE SCOPE & THANK YOU
# ═══════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "FUTURE SCOPE", font_size=14, color=PURPLE, bold=True)

# Future scope items
add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(3.8), BG_CARD, radius=0.04)
add_bullet_list(slide, Inches(1.0), Inches(1.5), Inches(5.0), Inches(3.2), [
    "🔄  Real-time JD scraping from LinkedIn & Naukri",
    "👥  Peer comparison within cohorts (anonymous)",
    "🎓  Alumni mentorship matching by target role",
    "📱  React Native mobile companion app",
    "🔗  LMS integration (Moodle/Canvas) for auto skill tracking",
    "🤝  Company-student direct messaging",
    "📈  Placement outcome prediction using ML",
], font_size=14, color=LIGHT_GRAY)

# Thank you
add_text_box(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(1),
             "Thank You!", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(0.6),
             "VidyaMitra — Your Knowledge Friend 🎓",
             font_size=20, color=CYAN, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(8.5), Inches(3.6), Inches(2.3), Inches(0.04), PURPLE)

add_text_box(slide, Inches(7), Inches(4.0), Inches(5.5), Inches(0.5),
             "Questions?",
             font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7), Inches(5.0), Inches(5.5), Inches(1.5),
             "Team: [Your Names]\n[College Name]\n[Contact / GitHub Link]",
             font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════

output_path = "/Users/srimannarayanadeevi/Documents/vidyamitra/VidyaMitra_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
